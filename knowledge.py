"""
SkyView Investment Advisors LLC
Claris Multi-Persona Platform — Knowledge Layer (RAG Pipeline)

Handles document ingestion, text chunking, embedding generation,
and similarity-based retrieval for grounding AI responses in
firm-approved content.
"""

import logging
import os
import re
from typing import Optional

from sqlalchemy import text as sql_text
from sqlalchemy.orm import Session as DBSession

from models import KnowledgeDocument, KnowledgeChunk

logger = logging.getLogger("skyview.knowledge")

# ─────────────────────────────────────────────────────────────────────────────
# CONFIGURATION
# ─────────────────────────────────────────────────────────────────────────────

EMBEDDING_MODEL = os.environ.get("EMBEDDING_MODEL", "voyage-3")
EMBEDDING_DIM = 1536
CHUNK_SIZE = 512        # tokens per chunk
CHUNK_OVERLAP = 50      # overlap between chunks
TOP_K = 5               # number of chunks to retrieve
SIMILARITY_THRESHOLD = 0.65  # minimum cosine similarity


# ─────────────────────────────────────────────────────────────────────────────
# TEXT EXTRACTION
# ─────────────────────────────────────────────────────────────────────────────

def extract_text_from_file(file_path: str, file_type: str) -> str:
    """
    Extract plain text from uploaded files.
    Supports PDF, DOCX, PPTX, TXT, CSV.
    """
    file_type = file_type.lower()

    if file_type in ("txt", "csv", "md", "json"):
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            return f.read()

    elif file_type == "pdf":
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(file_path)
            text = ""
            for page in doc:
                text += page.get_text() + "\n"
            doc.close()
            return text.strip()
        except ImportError:
            logger.error("PyMuPDF (fitz) not installed. Run: pip install PyMuPDF")
            raise

    elif file_type == "docx":
        try:
            from docx import Document
            doc = Document(file_path)
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except ImportError:
            logger.error("python-docx not installed. Run: pip install python-docx")
            raise

    elif file_type == "pptx":
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text)
            return "\n".join(texts)
        except ImportError:
            logger.error("python-pptx not installed. Run: pip install python-pptx")
            raise

    else:
        raise ValueError(f"Unsupported file type: {file_type}")


# ─────────────────────────────────────────────────────────────────────────────
# TEXT CHUNKING
# ─────────────────────────────────────────────────────────────────────────────

def chunk_text(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[str]:
    """
    Split text into overlapping chunks based on approximate token count.
    Uses a simple heuristic: 1 token ≈ 4 characters.
    """
    # Clean up whitespace
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = text.strip()

    if not text:
        return []

    char_chunk = chunk_size * 4
    char_overlap = overlap * 4

    chunks = []
    start = 0

    while start < len(text):
        end = start + char_chunk

        # Try to break at a paragraph or sentence boundary
        if end < len(text):
            # Look for paragraph break
            para_break = text.rfind("\n\n", start + char_chunk // 2, end + 200)
            if para_break > start:
                end = para_break

            # Fallback: sentence boundary
            elif text.rfind(". ", start + char_chunk // 2, end + 100) > start:
                end = text.rfind(". ", start + char_chunk // 2, end + 100) + 1

        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)

        start = end - char_overlap
        if start >= len(text):
            break

    logger.info(f"Chunked text into {len(chunks)} chunks (avg {sum(len(c) for c in chunks) // max(len(chunks), 1)} chars)")
    return chunks


# ─────────────────────────────────────────────────────────────────────────────
# EMBEDDING GENERATION
# ─────────────────────────────────────────────────────────────────────────────

def generate_embeddings(texts: list[str]) -> list[list[float]]:
    """
    Generate embeddings for a list of text strings.
    Supports Voyage AI (preferred) or OpenAI as fallback.
    """
    if not texts:
        return []

    voyage_key = os.environ.get("VOYAGE_API_KEY")
    openai_key = os.environ.get("OPENAI_API_KEY")

    if voyage_key:
        return _embed_voyage(texts, voyage_key)
    elif openai_key:
        return _embed_openai(texts, openai_key)
    else:
        logger.warning("No embedding API key found. Set VOYAGE_API_KEY or OPENAI_API_KEY.")
        return [[] for _ in texts]


def _embed_voyage(texts: list[str], api_key: str) -> list[list[float]]:
    """Generate embeddings using Voyage AI."""
    try:
        import voyageai
        client = voyageai.Client(api_key=api_key)
        result = client.embed(texts, model=EMBEDDING_MODEL, input_type="document")
        return result.embeddings
    except ImportError:
        logger.error("voyageai not installed. Run: pip install voyageai")
        raise


def _embed_openai(texts: list[str], api_key: str) -> list[list[float]]:
    """Generate embeddings using OpenAI (fallback)."""
    try:
        from openai import OpenAI
        client = OpenAI(api_key=api_key)
        response = client.embeddings.create(input=texts, model="text-embedding-ada-002")
        return [item.embedding for item in response.data]
    except ImportError:
        logger.error("openai not installed. Run: pip install openai")
        raise


def embed_query(query: str) -> list[float]:
    """Embed a single query string for retrieval."""
    voyage_key = os.environ.get("VOYAGE_API_KEY")
    openai_key = os.environ.get("OPENAI_API_KEY")

    if voyage_key:
        try:
            import voyageai
            client = voyageai.Client(api_key=voyage_key)
            result = client.embed([query], model=EMBEDDING_MODEL, input_type="query")
            return result.embeddings[0]
        except ImportError:
            pass

    if openai_key:
        try:
            from openai import OpenAI
            client = OpenAI(api_key=openai_key)
            response = client.embeddings.create(input=[query], model="text-embedding-ada-002")
            return response.data[0].embedding
        except ImportError:
            pass

    return []


# ─────────────────────────────────────────────────────────────────────────────
# DOCUMENT INGESTION
# ─────────────────────────────────────────────────────────────────────────────

def ingest_document(
    db: DBSession,
    title: str,
    content_text: str,
    category: str = "general",
    source_file: str = None,
    uploaded_by: str = None,
) -> KnowledgeDocument:
    """
    Ingest a document: store metadata, chunk text, generate embeddings,
    and save chunks with vectors to the database.
    """
    # Create document record
    doc = KnowledgeDocument(
        title=title,
        category=category,
        content_text=content_text,
        source_file=source_file,
        file_size_bytes=len(content_text.encode("utf-8")),
        uploaded_by=uploaded_by,
        status="pending",  # Requires compliance approval
    )
    db.add(doc)
    db.flush()  # Get the ID

    # Chunk the text
    chunks = chunk_text(content_text)

    if not chunks:
        logger.warning(f"No chunks generated for document '{title}'")
        db.commit()
        return doc

    # Generate embeddings
    embeddings = generate_embeddings(chunks)

    # Store chunks with embeddings
    for i, (chunk, embedding) in enumerate(zip(chunks, embeddings)):
        token_count = len(chunk) // 4  # Approximate
        db_chunk = KnowledgeChunk(
            document_id=doc.id,
            chunk_index=i,
            chunk_text=chunk,
            token_count=token_count,
            embedding_json=embedding if embedding else None,
        )
        db.add(db_chunk)

    db.commit()
    logger.info(f"Ingested '{title}': {len(chunks)} chunks with embeddings")
    return doc


# ─────────────────────────────────────────────────────────────────────────────
# RETRIEVAL
# ─────────────────────────────────────────────────────────────────────────────

class KnowledgeRetriever:
    """
    Retrieves relevant knowledge chunks using pgvector cosine similarity.
    Optionally filters by document category for persona-relevant results.
    """

    def __init__(self, db: DBSession, category_filter: list[str] = None):
        self.db = db
        self.category_filter = category_filter

    def retrieve(self, query: str, top_k: int = TOP_K) -> str:
        """
        Retrieve the most relevant knowledge chunks for a query.
        Returns formatted text ready for injection into the system prompt.
        """
        query_embedding = embed_query(query)

        if not query_embedding:
            logger.warning("No embedding generated for query — skipping RAG")
            return ""

        # Build the similarity search query
        embedding_str = f"[{','.join(str(x) for x in query_embedding)}]"

        sql = f"""
            SELECT
                kc.chunk_text,
                kd.title AS doc_title,
                kd.category,
                1 - (kc.embedding <=> '{embedding_str}'::vector) AS similarity
            FROM knowledge_chunks kc
            JOIN knowledge_documents kd ON kd.id = kc.document_id
            WHERE kd.is_active = true
              AND kd.status = 'approved'
        """

        if self.category_filter:
            cats = ",".join(f"'{c}'" for c in self.category_filter)
            sql += f" AND kd.category IN ({cats})"

        sql += f"""
            ORDER BY kc.embedding <=> '{embedding_str}'::vector
            LIMIT {top_k}
        """

        try:
            results = self.db.execute(sql_text(sql)).fetchall()
        except Exception as e:
            logger.error(f"RAG retrieval error: {e}")
            return ""

        if not results:
            return ""

        # Format retrieved chunks
        formatted = []
        for row in results:
            chunk_text, doc_title, category, similarity = row
            if similarity < SIMILARITY_THRESHOLD:
                continue
            formatted.append(
                f"[Source: {doc_title} ({category})] — Relevance: {similarity:.2f}\n"
                f"{chunk_text}"
            )

        if not formatted:
            return ""

        context = "\n\n---\n\n".join(formatted)
        logger.info(f"RAG retrieved {len(formatted)} relevant chunks for query")
        return context
